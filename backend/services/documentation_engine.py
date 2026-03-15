import io
import json
import logging
import os
import re
import time
from datetime import datetime
from pathlib import Path
from typing import AsyncGenerator, Generator, List

import requests
from fastapi import HTTPException, UploadFile, status
from pypdf import PdfReader

# Logger setup
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s - %(message)s",
)
logger = logging.getLogger("documentation_engine")


class DocumentationEngine:
    def __init__(self) -> None:
        self.base_url = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434").rstrip("/")
        self.model = os.getenv("OLLAMA_DOC_MODEL", "llama3:latest")
        self.generate_url = f"{self.base_url}/api/generate"
        self.max_retries = int(os.getenv("OLLAMA_DOC_MAX_RETRIES", "2"))
        self.retry_backoff_seconds = float(os.getenv("OLLAMA_DOC_RETRY_BACKOFF_SECONDS", "1.5"))
        logger.info(
            "DocumentationEngine initialized",
            extra={"base_url": self.base_url, "model": self.model},
        )

    def _extract_text_from_pdf(self, content: bytes) -> str:
        try:
            reader = PdfReader(io.BytesIO(content))
            chunks: List[str] = []
            for page in reader.pages:
                text = page.extract_text() or ""
                if text.strip():
                    chunks.append(text.strip())
            return "\n\n".join(chunks)
        except Exception as exc:
            logger.error("Failed to parse PDF", extra={"error": str(exc)})
            return ""

    def _extract_text_from_pptx(self, content: bytes) -> str:
        try:
            from pptx import Presentation
        except Exception as exc:
            logger.error("python-pptx import failed", extra={"error": str(exc)})
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail="python-pptx is not installed. Please install dependencies.",
            ) from exc

        try:
            presentation = Presentation(io.BytesIO(content))
            lines: List[str] = []
            for slide_index, slide in enumerate(presentation.slides, start=1):
                slide_lines: List[str] = []
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text and text.strip():
                        slide_lines.append(text.strip())
                if slide_lines:
                    lines.append(f"Slide {slide_index}:\n" + "\n".join(slide_lines))
            return "\n\n".join(lines)
        except Exception as exc:
            logger.error("Failed to parse PPTX", extra={"error": str(exc)})
            return ""

    def _extract_text_from_file(self, file_name: str, content: bytes) -> str:
        lowered = file_name.lower()

        if lowered.endswith((".txt", ".md", ".csv", ".json")):
            for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin1"):
                try:
                    return content.decode(encoding)
                except UnicodeDecodeError:
                    continue
            return ""

        if lowered.endswith(".pdf"):
            return self._extract_text_from_pdf(content)

        if lowered.endswith(".pptx"):
            return self._extract_text_from_pptx(content)

        if lowered.endswith(".ppt"):
            return "Unsupported binary .ppt format. Please upload .pptx for text extraction."

        return ""

    async def _build_documents_context(self, files: List[UploadFile]) -> str:
        if not files:
            return ""

        sections: List[str] = []
        for upload in files:
            content = await upload.read()
            extracted_text = self._extract_text_from_file(upload.filename or "document", content)
            if extracted_text.strip():
                sections.append(
                    f"### SOURCE FILE: {upload.filename}\n{extracted_text[:50000]}"
                )

        return "\n\n".join(sections)

    def _build_prompt(self, user_prompt: str, context_text: str) -> str:
        context_block = context_text.strip() or "No attachment text provided."
        return f"""
You are a documentation engine.
Generate a polished, editable documentation draft in markdown.

Rules:
1. Use clear headings, tables where useful, and concise sections.
2. If source context is present, cite source file names inline as plain text references.
3. Include: Executive Summary, Key Points, Detailed Sections, Risks/Notes, Next Steps.
4. Keep output professional and complete.
5. Return plain markdown text only.

USER PROMPT:
{user_prompt}

SOURCE CONTENT:
{context_block}
""".strip()

    def _should_retry(self, exc: requests.RequestException) -> bool:
        response = getattr(exc, "response", None)
        if response is None:
            return True
        status_code = int(getattr(response, "status_code", 0))
        return status_code >= 500 or status_code == 429

    def stream_document(self, prompt: str) -> Generator[str, None, None]:
        payload = {
            "model": self.model,
            "prompt": prompt,
            "stream": True,
        }

        for attempt in range(self.max_retries + 1):
            try:
                with requests.post(self.generate_url, json=payload, stream=True, timeout=900) as response:
                    response.raise_for_status()
                    for raw_line in response.iter_lines(decode_unicode=True):
                        if not raw_line:
                            continue
                        try:
                            row = json.loads(raw_line)
                        except Exception:
                            continue
                        chunk = str(row.get("response", ""))
                        if chunk:
                            yield chunk
                        if bool(row.get("done", False)):
                            return
                return
            except requests.RequestException as exc:
                is_last_attempt = attempt >= self.max_retries
                if is_last_attempt or not self._should_retry(exc):
                    logger.error("Failed while streaming documentation", extra={"error": str(exc)})
                    raise HTTPException(
                        status_code=status.HTTP_502_BAD_GATEWAY,
                        detail=f"Failed while streaming documentation: {str(exc)}",
                    ) from exc
                sleep_seconds = self.retry_backoff_seconds * (attempt + 1)
                logger.warning(
                    "Transient documentation model error, retrying",
                    extra={"attempt": attempt + 1, "sleep_seconds": sleep_seconds, "error": str(exc)},
                )
                time.sleep(sleep_seconds)

    async def stream_documentation_from_inputs(
        self,
        *,
        user_prompt: str,
        files: List[UploadFile],
    ) -> AsyncGenerator[str, None]:
        context_text = await self._build_documents_context(files)
        final_prompt = self._build_prompt(user_prompt, context_text)
        for chunk in self.stream_document(final_prompt):
            yield chunk

    def save_document(self, title: str, content: str) -> tuple[str, int]:
        safe_title = re.sub(r"[^a-zA-Z0-9_-]+", "-", title.strip().lower()) or "documentation"
        output_dir = Path("generated_documents")
        output_dir.mkdir(parents=True, exist_ok=True)
        timestamp = datetime.utcnow().strftime("%Y%m%d-%H%M%S")
        output_path = output_dir / f"{safe_title}-{timestamp}.md"
        output_path.write_text(content, encoding="utf-8")
        return str(output_path), len(content.encode("utf-8"))
